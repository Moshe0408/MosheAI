"""כלי יצירת קבצים עבור MosheAI - מותאם ל-Vercel"""

import json
import base64
import datetime
import traceback
from pathlib import Path

# Vercel: only /tmp is writable
OUTPUT_DIR = Path("/tmp/mosheai_outputs")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# ── ייבוא חבילות ─────────────────────────────

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    PPTX_OK = True
except ImportError:
    PPTX_OK = False

try:
    from docx import Document
    from docx.shared import Pt as DPt, RGBColor as DRGB
    DOCX_OK = True
except ImportError:
    DOCX_OK = False

try:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    MATPLOTLIB_OK = True
except ImportError:
    MATPLOTLIB_OK = False

# ── סכמת כלים ─────────────────────────────────

TOOLS_SCHEMA_GROQ = [
    {
        "type": "function",
        "function": {
            "name": "create_presentation",
            "description": "יוצר מצגת PowerPoint מקצועית עם שקופיות, כותרות ונקודות.",
            "parameters": {
                "type": "object",
                "properties": {
                    "title":       {"type": "string"},
                    "filename":    {"type": "string", "description": "שם קובץ ללא סיומת"},
                    "theme_color": {"type": "string", "description": "צבע hex ראשי, ברירת מחדל #1F4E79"},
                    "slides": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "heading": {"type": "string"},
                                "bullets": {"type": "array", "items": {"type": "string"}},
                                "notes":   {"type": "string"}
                            },
                            "required": ["heading", "bullets"]
                        }
                    }
                },
                "required": ["title", "filename", "slides"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "create_word_report",
            "description": "יוצר דוח Word מקצועי עם כותרות, פסקאות וטבלאות.",
            "parameters": {
                "type": "object",
                "properties": {
                    "title":    {"type": "string"},
                    "filename": {"type": "string"},
                    "sections": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "heading": {"type": "string"},
                                "content": {"type": "string"},
                                "table": {
                                    "type": "object",
                                    "properties": {
                                        "headers": {"type": "array", "items": {"type": "string"}},
                                        "rows":    {"type": "array", "items": {"type": "array", "items": {"type": "string"}}}
                                    }
                                }
                            },
                            "required": ["heading", "content"]
                        }
                    }
                },
                "required": ["title", "filename", "sections"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "create_chart",
            "description": "יוצר גרף סטטיסטי (עמודות/עוגה/קו) ושומר PNG.",
            "parameters": {
                "type": "object",
                "properties": {
                    "chart_type": {"type": "string", "enum": ["bar", "pie", "line", "barh"]},
                    "title":      {"type": "string"},
                    "filename":   {"type": "string"},
                    "labels":     {"type": "array", "items": {"type": "string"}},
                    "values":     {"type": "array", "items": {"type": "number"}},
                    "xlabel":     {"type": "string"},
                    "ylabel":     {"type": "string"},
                    "colors":     {"type": "array", "items": {"type": "string"}}
                },
                "required": ["chart_type", "title", "filename", "labels", "values"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "recall_memory",
            "description": "מחזיר את ההיסטוריה, הטעויות הקודמות והסטטיסטיקה של הסוכן.",
            "parameters": {"type": "object", "properties": {}, "required": []}
        }
    }
]


# ── מימוש כלים ────────────────────────────────

def run_tool(name: str, args: dict, memory: dict) -> dict:
    try:
        if name == "create_presentation":
            return _pptx(args)
        elif name == "create_word_report":
            return _docx(args)
        elif name == "create_chart":
            return _chart(args)
        elif name == "recall_memory":
            return _recall(memory)
        else:
            return {"error": f"כלי לא מוכר: {name}"}
    except Exception as e:
        return {"error": str(e), "traceback": traceback.format_exc()[-400:]}


def _hex(h: str) -> tuple:
    h = h.lstrip("#")
    try:
        return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    except Exception:
        return 31, 78, 121


def _pptx(args: dict) -> dict:
    if not PPTX_OK:
        return {"error": "python-pptx לא מותקן"}

    r, g, b = _hex(args.get("theme_color", "#1F4E79"))
    filename = args["filename"].removesuffix(".pptx") + ".pptx"
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    sl = prs.slides.add_slide(prs.slide_layouts[0])
    sl.shapes.title.text = args["title"]
    sl.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(r, g, b)
    sl.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(36)
    try:
        sl.placeholders[1].text = datetime.datetime.now().strftime("%d/%m/%Y")
    except Exception:
        pass

    for s in args["slides"]:
        sl = prs.slides.add_slide(prs.slide_layouts[1])
        sl.shapes.title.text = s["heading"]
        sl.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(r, g, b)
        sl.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
        tf = sl.placeholders[1].text_frame
        tf.clear()
        for i, bullet in enumerate(s["bullets"]):
            p = tf.add_paragraph() if i else tf.paragraphs[0]
            p.text = bullet
            p.font.size = Pt(18)
        if s.get("notes"):
            sl.notes_slide.notes_text_frame.text = s["notes"]

    path = OUTPUT_DIR / filename
    prs.save(str(path))

    # Base64 for direct download (Vercel serverless safe)
    data_b64 = base64.b64encode(path.read_bytes()).decode()

    return {
        "success":  True,
        "path":     str(path),
        "filename": filename,
        "type":     "pptx",
        "slides":   len(args["slides"]) + 1,
        "data_b64": data_b64,
        "mime":     "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    }


def _docx(args: dict) -> dict:
    if not DOCX_OK:
        return {"error": "python-docx לא מותקן"}

    filename = args["filename"].removesuffix(".docx") + ".docx"
    doc = Document()
    h = doc.add_heading(args["title"], level=0)
    h.runs[0].font.color.rgb = DRGB(0x1F, 0x4E, 0x79)
    doc.add_paragraph(f"נוצר: {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}")
    doc.add_paragraph("")

    for sec in args["sections"]:
        doc.add_heading(sec["heading"], level=1)
        doc.add_paragraph(sec["content"])
        if tbl := sec.get("table"):
            headers = tbl.get("headers", [])
            rows    = tbl.get("rows", [])
            if headers:
                table = doc.add_table(1 + len(rows), len(headers))
                table.style = "Light Shading Accent 1"
                for ci, h_text in enumerate(headers):
                    table.rows[0].cells[ci].text = h_text
                for ri, row in enumerate(rows):
                    for ci, val in enumerate(row):
                        table.rows[ri + 1].cells[ci].text = str(val)
        doc.add_paragraph("")

    path = OUTPUT_DIR / filename
    doc.save(str(path))

    data_b64 = base64.b64encode(path.read_bytes()).decode()

    return {
        "success":  True,
        "path":     str(path),
        "filename": filename,
        "type":     "docx",
        "data_b64": data_b64,
        "mime":     "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    }


def _chart(args: dict) -> dict:
    if not MATPLOTLIB_OK:
        return {"error": "matplotlib לא מותקן"}

    filename = args["filename"].removesuffix(".png") + ".png"
    labels   = args["labels"]
    values   = args["values"]
    colors   = args.get("colors") or ["#7c3aed","#2563eb","#10b981","#f59e0b","#ef4444","#06b6d4","#8b5cf6"]

    fig, ax = plt.subplots(figsize=(10, 6))
    fig.patch.set_facecolor("#111827")
    ax.set_facecolor("#111827")
    ax.tick_params(colors="#e2e8f0")
    ax.xaxis.label.set_color("#94a3b8")
    ax.yaxis.label.set_color("#94a3b8")
    for spine in ax.spines.values():
        spine.set_edgecolor("#374151")

    ct = args["chart_type"]
    c  = (colors * max(len(labels), 1))[:len(labels)]

    if ct == "bar":
        bars = ax.bar(labels, values, color=c, edgecolor="#1f2937", linewidth=0.5)
        ax.bar_label(bars, fmt="%.1f", padding=3, color="#e2e8f0", fontsize=9)
        ax.set_xlabel(args.get("xlabel",""), color="#94a3b8")
        ax.set_ylabel(args.get("ylabel",""), color="#94a3b8")
        ax.tick_params(axis="x", rotation=30)
    elif ct == "barh":
        bars = ax.barh(labels, values, color=c, edgecolor="#1f2937", linewidth=0.5)
        ax.bar_label(bars, fmt="%.1f", padding=3, color="#e2e8f0", fontsize=9)
        ax.set_xlabel(args.get("xlabel",""), color="#94a3b8")
    elif ct == "pie":
        ax.pie(values, labels=labels, colors=c, autopct="%1.1f%%", startangle=140,
               wedgeprops={"edgecolor":"#111827","linewidth":1.5},
               textprops={"color":"#e2e8f0"})
        ax.axis("equal")
    elif ct == "line":
        ax.plot(labels, values, marker="o", color=colors[0], linewidth=2.5,
                markersize=8, markerfacecolor="#111827", markeredgewidth=2)
        ax.fill_between(range(len(labels)), values, alpha=0.15, color=colors[0])
        ax.set_xlabel(args.get("xlabel",""), color="#94a3b8")
        ax.set_ylabel(args.get("ylabel",""), color="#94a3b8")
        ax.set_xticks(range(len(labels)))
        ax.set_xticklabels(labels, rotation=30, ha="right")

    ax.set_title(args["title"], color="#f1f5f9", fontsize=14, fontweight="bold", pad=16)
    plt.tight_layout()

    path = OUTPUT_DIR / filename
    plt.savefig(str(path), dpi=150, bbox_inches="tight", facecolor="#111827")
    plt.close(fig)

    thumb_b64 = base64.b64encode(path.read_bytes()).decode()

    return {
        "success":    True,
        "path":       str(path),
        "filename":   filename,
        "type":       "png",
        "thumbnail":  thumb_b64,
        "data_b64":   thumb_b64,
        "mime":       "image/png"
    }


def _recall(memory: dict) -> dict:
    s = memory.get("stats", {})
    return {
        "stats":           s,
        "recent_errors":   memory.get("errors",   [])[-5:],
        "recent_sessions": memory.get("sessions", [])[-5:]
    }


def list_outputs() -> list:
    files = []
    try:
        for p in sorted(OUTPUT_DIR.iterdir(), key=lambda x: -x.stat().st_mtime):
            if p.suffix in (".pptx", ".docx", ".png", ".xlsx"):
                files.append({
                    "filename": p.name,
                    "size":     _human_size(p.stat().st_size),
                    "modified": datetime.datetime.fromtimestamp(p.stat().st_mtime).strftime("%d/%m/%Y %H:%M"),
                    "type":     p.suffix[1:]
                })
    except Exception:
        pass
    return files


def _human_size(b: int) -> str:
    for unit in ("B", "KB", "MB"):
        if b < 1024:
            return f"{b:.0f} {unit}"
        b /= 1024
    return f"{b:.1f} GB"
